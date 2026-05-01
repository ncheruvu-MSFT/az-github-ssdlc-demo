"""Generate a 2-slide PowerPoint: Problem Statement & Solution Architecture."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_RED = RGBColor(0xF8, 0x51, 0x49)
ACCENT_ORANGE = RGBColor(0xF0, 0x88, 0x3E)
ACCENT_PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)
MID_GRAY = RGBColor(0x8B, 0x94, 0x9E)
CARD_BG = RGBColor(0x16, 0x1B, 0x22)
BORDER_GRAY = RGBColor(0x30, 0x36, 0x3D)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape


def set_text(shape, text, font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=12, color=WHITE, bold=False, space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_bullet_card(slide, left, top, width, height, title, bullets, icon_color, border_color=None):
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG, border_color or icon_color)
    card.text_frame.word_wrap = True
    card.text_frame.margin_left = Pt(14)
    card.text_frame.margin_right = Pt(10)
    card.text_frame.margin_top = Pt(10)
    card.text_frame.margin_bottom = Pt(8)
    tf = set_text(card, title, font_size=14, color=icon_color, bold=True)
    for b in bullets:
        add_para(tf, f"  •  {b}", font_size=11, color=LIGHT_GRAY, space_before=Pt(3))
    return card


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

# Title
tx = slide1.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
set_text(tx, "The Problem: Why Traditional SDLC Falls Short", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_para(tx.text_frame, "Enterprise teams face compounding security & delivery challenges that manual processes can't solve",
         font_size=14, color=MID_GRAY, space_before=Pt(6))

# ── Row 1 — 3 cards ──
card_w = Inches(3.8)
card_h = Inches(2.2)
gap = Inches(0.45)
row1_top = Inches(1.65)
row1_left = Inches(0.6)

add_bullet_card(slide1, row1_left, row1_top, card_w, card_h,
    "🔓  Security as an Afterthought",
    ["Vulnerabilities found late in production",
     "Secret sprawl — keys in code & configs",
     "No automated SAST/DAST in CI pipelines",
     "OWASP Top 10 gaps undetected until audit"],
    ACCENT_RED)

add_bullet_card(slide1, row1_left + card_w + gap, row1_top, card_w, card_h,
    "🐢  Slow, Manual Release Process",
    ["Weeks-long deploy cycles with manual gates",
     "No immutable artifacts — 'it works on my machine'",
     "Credential rotation is manual & error-prone",
     "Environment drift between dev/staging/prod"],
    ACCENT_ORANGE)

add_bullet_card(slide1, row1_left + 2*(card_w + gap), row1_top, card_w, card_h,
    "🔍  Zero Traceability & Visibility",
    ["No link from requirement → code → deployment",
     "Audit readiness takes weeks of manual effort",
     "Dependency risks invisible until CVE hits prod",
     "Siloed tools with no unified security dashboard"],
    ACCENT_PURPLE)

# ── Row 2 — 3 cards ──
row2_top = Inches(4.15)

add_bullet_card(slide1, row1_left, row2_top, card_w, card_h,
    "⚙️  Infrastructure Inconsistency",
    ["Hand-crafted environments = configuration drift",
     "No IaC scanning — misconfigs reach production",
     "Public endpoints exposed by default",
     "Missing diagnostics & monitoring blind spots"],
    ACCENT_BLUE)

add_bullet_card(slide1, row1_left + card_w + gap, row2_top, card_w, card_h,
    "📦  Supply Chain Risks",
    ["No dependency review on pull requests",
     "Unsigned container images in production",
     "Typosquatting & malicious packages undetected",
     "No SBOM or provenance verification"],
    ACCENT_GREEN)

add_bullet_card(slide1, row1_left + 2*(card_w + gap), row2_top, card_w, card_h,
    "🧩  Developer Friction",
    ["Context-switching across disconnected tools",
     "Security findings without remediation guidance",
     "No AI-assisted development workflow",
     "Manual work-item tracking & status updates"],
    RGBColor(0xE0, 0x70, 0x70))

# Bottom bar
bottom = slide1.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12), Inches(0.5))
set_text(bottom,
    "Result:  Slower releases  •  Higher risk  •  Audit failures  •  Developer burnout  •  No end-to-end traceability",
    font_size=13, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide2, BG_DARK)

# Title
tx2 = slide2.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.65))
set_text(tx2, "Solution: Modern Secure SDLC — GitHub + Azure Enterprise Pipeline",
         font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tx2.text_frame,
    "End-to-end DevSecOps  •  Zero stored credentials  •  Progressive deployment  •  Full traceability",
    font_size=12, color=MID_GRAY, space_before=Pt(4), alignment=PP_ALIGN.CENTER)

# ── Layer layout — 6 horizontal bands ──
layer_left = Inches(0.35)
layer_w = Inches(12.6)
band_gap = Pt(6)

def add_layer_band(slide, top, height, label, label_color, items):
    """Add a horizontal layer band with label + item cards."""
    # Label column
    lbl = add_rounded_rect(slide, layer_left, top, Inches(1.9), height, CARD_BG, label_color)
    lbl.text_frame.word_wrap = True
    lbl.text_frame.margin_left = Pt(8)
    lbl.text_frame.margin_top = Pt(6)
    set_text(lbl, label, font_size=11, color=label_color, bold=True, alignment=PP_ALIGN.CENTER)
    lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Item cards
    items_left = layer_left + Inches(2.05)
    items_total_w = layer_w - Inches(2.05)
    card_gap = Pt(6)
    n = len(items)
    cw = (items_total_w - Emu(card_gap * (n - 1))) / n if n else items_total_w

    for i, (title, desc) in enumerate(items):
        cl = items_left + (cw + Emu(card_gap)) * i
        card = add_rounded_rect(slide, cl, top, cw, height, CARD_BG, BORDER_GRAY)
        card.text_frame.word_wrap = True
        card.text_frame.margin_left = Pt(6)
        card.text_frame.margin_right = Pt(4)
        card.text_frame.margin_top = Pt(4)
        card.text_frame.margin_bottom = Pt(2)
        tf = set_text(card, title, font_size=10, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        if desc:
            add_para(tf, desc, font_size=8, color=LIGHT_GRAY, space_before=Pt(2), alignment=PP_ALIGN.CENTER)


band_h = Inches(0.82)
y = Inches(1.15)
spacing = Inches(0.88)

# Layer 1 — Work Management
add_layer_band(slide2, y, band_h, "① WORK\nMANAGEMENT", RGBColor(0x79, 0xC0, 0xFF), [
    ("ADO Boards", "Epics → Stories → Tasks\nSprint planning & velocity"),
    ("GitHub Issues", "Labels, milestones\nProjects v2 boards"),
    ("AB# Integration", "Bidirectional links\nCommit → Work item"),
    ("Copilot Agent Bridge", "ADO → repository_dispatch\nAI-assisted PR creation"),
])

# Layer 2 — GitHub Platform + GHAS
y += spacing
add_layer_band(slide2, y, band_h, "② GITHUB\nPLATFORM + GHAS", ACCENT_GREEN, [
    ("Branch Protection", "2 approvals, signed commits\nCODEOWNERS enforcement"),
    ("CodeQL SAST", "C# & Python analysis\nSecurity + quality queries"),
    ("Secret Scanning", "Push protection\n200+ secret patterns"),
    ("Dependabot", "NuGet, pip, Actions\nAuto-PR for updates"),
    ("Dep Review + Malware", "Block vulnerable deps\nTyposquatting detection"),
])

# Layer 3 — CI Pipeline
y += spacing
add_layer_band(slide2, y, band_h, "③ CI PIPELINE\n(every PR)", ACCENT_ORANGE, [
    (".NET Build + Test", "xUnit, FluentAssertions\n≥60% coverage gate"),
    ("Python Lint + Test", "Ruff + pytest + httpx\nBandit SAST"),
    ("Container Scan", "MS Defender\nImage vuln scan"),
    ("IaC Scan", "Checkov\nBicep misconfig check"),
    ("Coverage + Review", "PR comments\nDep review gate"),
])

# Layer 4 — CD Pipeline
y += spacing
add_layer_band(slide2, y, band_h, "④ CD PIPELINE\n(on merge)", ACCENT_RED, [
    ("OIDC Auth", "Zero stored credentials\nFederated identity"),
    ("Notation Signing", "Notary v2 image signing\nSupply chain integrity"),
    ("DEV → STAGING", "Auto-deploy\nHealth checks"),
    ("Staging Gate Report", "Go/no-go readiness\nCI + signature verify"),
    ("Manual Gate → PROD", "Reviewer approval\nTeams notification"),
])

# Layer 5 — Azure Platform
y += spacing
add_layer_band(slide2, y, band_h, "⑤ AZURE\nPLATFORM", ACCENT_PURPLE, [
    ("Functions (.NET 8)", "Isolated worker\nDurable orchestration"),
    ("Container Apps", "C# Minimal API\nPython FastAPI"),
    ("Service Bus Premium", "Queues + Topics\nDead-letter & audit"),
    ("Key Vault + RBAC", "Managed Identity\nPrivate endpoints"),
    ("VNet + NSG", "Subnet segmentation\nDeny-all default"),
])

# Layer 6 — IaC + Monitoring
y += spacing
add_layer_band(slide2, y, band_h, "⑥ IaC +\nMONITORING", RGBColor(0x56, 0xD3, 0x64), [
    ("Bicep (AVM)", "Subscription-scope\nEnv parameterization"),
    ("App Insights", "Distributed tracing\nPerformance metrics"),
    ("Log Analytics", "Centralized logs\nDiagnostic settings"),
    ("Alert Rules", "Action groups\nProactive notifications"),
])

# ── Flow arrows (text-based) ──
arrow = slide2.shapes.add_textbox(Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.35))
tf_arrow = set_text(arrow,
    "Flow:  Work Item  →  Feature Branch  →  PR + CI (8 checks)  →  Merge  →  Build + Sign  →  DEV  →  STAGING  →  Gate  →  PROD  →  Telemetry",
    font_size=11, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ── Save ──
out_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "Modern-SSDLC-GitHub-Azure.pptx")
prs.save(out_path)
print(f"✅ Saved: {out_path}")
